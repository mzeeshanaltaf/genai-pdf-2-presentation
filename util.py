import pdfplumber
from pptx import Presentation
from io import BytesIO
from langchain_core.prompts import ChatPromptTemplate
from schemas import *


# Function to extract text from PDF file
def extract_text_from_pdf(pdf_file_path):
    # Open the PDF file using pdfplumber
    with pdfplumber.open(pdf_file_path) as pdf:
        # Initialize an empty string to store extracted text
        extracted_text = ""

        # Loop through all the pages and extract text
        for page in pdf.pages:
            extracted_text += page.extract_text()

    return extracted_text


# This function uses an LLM to create well-structured slides with key points based on provided text
def generate_presentation(number_of_slides, number_of_bullet_points, extracted_text, llm):
    # Create prompt
    prompt_template = ChatPromptTemplate.from_template(PROMPT_TEMPLATE_SLIDES)
    prompt = prompt_template.format(number_of_slides=number_of_slides,
                                    number_of_bullet_points=number_of_bullet_points,
                                    text=extracted_text)

    # Get structured output from LLM
    structured_llm = llm.with_structured_output(CreatePresentation)
    structured_response = structured_llm.invoke(prompt)
    presentation_data = structured_response.dict()
    return presentation_data


# This function uses an LLM to create podcast script based on provided text
def generate_podcast(extracted_text, llm):
    prompt_template = ChatPromptTemplate.from_template(PROMPT_TEMPLATE_PODCAST)
    prompt_podcast = prompt_template.format(text=extracted_text)

    # Get structured output from llm
    structured_llm = llm.with_structured_output(CreatePodcast)
    structured_response = structured_llm.invoke(prompt_podcast)
    podcast_data = structured_response.dict()
    return podcast_data


# This function generates PowerPoint presentation (pptx) from the presentation content provided to it
def text_to_presentation(slide_contents, presentation_title, slide_notes):
    # Split the text by slides (assuming each slide's content is separated by two newlines)
    slides_content = slide_contents.strip().split("\n\n")

    # Create a new PowerPoint presentation object
    prs = Presentation()

    # Create the first slide with the presentation title
    title_slide_layout = prs.slide_layouts[0]  # Use slide layout 0 (Title Slide)
    title_slide = prs.slides.add_slide(title_slide_layout)
    title_slide.shapes.title.text = presentation_title

    # Add notes to the first slide
    notes_slide = title_slide.notes_slide
    notes_slide.notes_text_frame.text = slide_notes[0]

    # Create the remaining slides with content
    for idx, slide_content in enumerate(slides_content):
        # Split the content into title and bullet points
        lines = slide_content.split("\n")
        slide_title = lines[0]  # First line is the title
        bullet_points = lines[1:]  # Remaining lines are bullet points

        # Add a new slide (using the layout with a title and content)
        slide_layout = prs.slide_layouts[1]  # Use slide layout 1 (Title and Content)
        slide = prs.slides.add_slide(slide_layout)

        # Set the title for the slide
        slide.shapes.title.text = slide_title

        # Set the bullet points for the slide
        content_shape = slide.shapes.placeholders[1]
        text_frame = content_shape.text_frame

        for point in bullet_points:
            p = text_frame.add_paragraph()
            p.text = point

        # Add notes to the slide
        if idx + 1 < len(slide_notes):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_notes[idx + 1]

    # Save the presentation to a BytesIO object instead of a file
    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)  # Move the cursor to the start of the BytesIO object
    return pptx_io


# This function extracts presentation data from structured LLM response
def extract_presentation_data(presentation_data):
    # Extract presentation title
    presentation_title = presentation_data['presentation_title']

    # Extract slide content
    slide_contents = ""
    for slide in presentation_data['slides']:
        slide_contents += f"{slide['slide_title']}\n"
        for bullet in slide['bullet_points']:
            slide_contents += f"    {bullet}\n"
        slide_contents += "\n"  # For separating each slide

    # Extract slide notes
    slide_notes = ["This is the title slide."]  # Default first note
    for slide in presentation_data['slides']:
        slide_notes.append(f"'{slide['slide_notes']}'")

    return presentation_title, slide_contents, slide_notes
